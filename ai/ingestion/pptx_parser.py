from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

from .base_parser import BaseParser, ParsedDocument

SUPPORTED_EXTENSIONS = {".pptx"}


class PPTXParser(BaseParser):
    """
    Parses PowerPoint files using python-pptx.

    Each slide becomes one "page" in the ParsedDocument.
    Text is extracted in reading order: title → body → speaker notes.

    Key insight: PPTX slides often have very little text per slide.
    The chunker should treat multi-slide sequences as one chunk rather
    than chunking each slide individually — we handle this with metadata.
    """

    def supports(self, file_path: str | Path) -> bool:
        return Path(file_path).suffix.lower() in SUPPORTED_EXTENSIONS

    def parse(self, file_path: str | Path) -> ParsedDocument:
        path = self._validate_file(file_path)

        if not self.supports(path):
            raise ValueError(f"PPTXParser does not support: {path.suffix}")

        prs = Presentation(str(path))
        pages = []
        all_text_parts = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_title = self._get_slide_title(slide)
            body_text = self._get_body_text(slide)
            notes_text = self._get_notes_text(slide)

            # Combine title + body + notes into one text block per slide
            parts = []
            if slide_title:
                parts.append(f"[Slide {slide_num}: {slide_title}]")
            if body_text:
                parts.append(body_text)
            if notes_text:
                parts.append(f"[Speaker notes: {notes_text}]")

            slide_text = "\n".join(parts).strip()

            if slide_text:
                pages.append({
                    "page_num": slide_num,
                    "text": slide_text,
                    "metadata": {
                        "title": slide_title,
                        "has_notes": bool(notes_text),
                        "shape_count": len(slide.shapes),
                    },
                })
                all_text_parts.append(slide_text)

        full_text = "\n\n".join(all_text_parts)

        # Presentation-level metadata
        props = prs.core_properties
        metadata = {
            "author": props.author or "",
            "title": props.title or "",
            "slide_count": len(prs.slides),
            "file_size_bytes": path.stat().st_size,
        }

        title = props.title or path.stem

        return ParsedDocument(
            source_path=str(path),
            file_type="pptx",
            title=title,
            full_text=full_text,
            pages=pages,
            metadata=metadata,
        )

    def _get_slide_title(self, slide) -> str:
        """Get the title placeholder text, if any."""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            return slide.shapes.title.text.strip()
        return ""

    def _get_body_text(self, slide) -> str:
        """
        Extract all non-title text from a slide.
        Iterates over all shapes, skips the title placeholder.
        Preserves bullet structure with indentation.
        """
        texts = []
        for shape in slide.shapes:
            # Skip the title shape (already handled)
            if shape == slide.shapes.title:
                continue
            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                # Preserve bullet hierarchy with indentation
                indent = "  " * (para.level or 0)
                texts.append(f"{indent}{text}")

        return "\n".join(texts)

    def _get_notes_text(self, slide) -> str:
        """Extract speaker notes text."""
        try:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            return tf.text.strip() if tf else ""
        except Exception:
            return ""