"""
Tests for the ingestion layer.
Creates real test files in /tmp, parses them, and validates the output.
Run with: python -m pytest tests/test_ingestion.py -v
"""

import os
import tempfile
import pytest
from pathlib import Path

from ai.ingestion import parse_document, get_supported_extensions, ParsedDocument


# ── Helpers ──────────────────────────────────────────────────────────────────

def make_md_file(content: str, suffix=".md") -> Path:
    tmp = tempfile.NamedTemporaryFile(suffix=suffix, mode="w", delete=False, encoding="utf-8")
    tmp.write(content)
    tmp.close()
    return Path(tmp.name)

def make_pdf_file() -> Path:
    """Create a minimal real PDF using PyMuPDF."""
    import fitz
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 100), "Introduction to RAG Systems", fontsize=18)
    page.insert_text((50, 140), "Retrieval-Augmented Generation combines search with LLMs.")
    page.insert_text((50, 165), "The system first retrieves relevant documents, then generates answers.")
    page2 = doc.new_page()
    page2.insert_text((50, 100), "How Embeddings Work", fontsize=18)
    page2.insert_text((50, 140), "Text is converted to vectors using embedding models.")
    tmp = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
    doc.save(tmp.name)
    doc.close()
    return Path(tmp.name)

def make_docx_file() -> Path:
    from docx import Document
    doc = Document()
    doc.add_heading("Research Summary", level=1)
    doc.add_paragraph("This document summarizes key findings from our research.")
    doc.add_heading("Methodology", level=2)
    doc.add_paragraph("We used a hybrid retrieval approach combining dense and sparse search.")
    # Add a table
    table = doc.add_table(rows=2, cols=3)
    table.cell(0, 0).text = "Method"
    table.cell(0, 1).text = "Precision"
    table.cell(0, 2).text = "Recall"
    table.cell(1, 0).text = "BM25"
    table.cell(1, 1).text = "0.72"
    table.cell(1, 2).text = "0.68"
    tmp = tempfile.NamedTemporaryFile(suffix=".docx", delete=False)
    doc.save(tmp.name)
    return Path(tmp.name)

def make_pptx_file() -> Path:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "RAG Architecture"
    slide.placeholders[1].text = "Key components:\n• Ingestion\n• Embedding\n• Retrieval\n• Generation"
    # Add speaker notes
    notes = slide.notes_slide.notes_text_frame
    notes.text = "Emphasize the retrieval quality as the bottleneck."
    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Results"
    slide2.placeholders[1].text = "Achieved 94% accuracy on benchmark dataset."
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    return Path(tmp.name)


# ── Markdown tests ────────────────────────────────────────────────────────────

class TestMarkdownParser:
    def test_basic_parse(self):
        path = make_md_file("# Hello World\n\nThis is a test document.\n\nIt has two paragraphs.")
        doc = parse_document(path)
        assert isinstance(doc, ParsedDocument)
        assert doc.file_type == "md"
        assert "Hello World" in doc.full_text
        assert "two paragraphs" in doc.full_text
        os.unlink(path)

    def test_title_extracted_from_h1(self):
        path = make_md_file("# My Research Paper\n\nContent here.")
        doc = parse_document(path)
        assert doc.title == "My Research Paper"
        os.unlink(path)

    def test_sections_split_by_headings(self):
        content = "# Title\n\nIntro.\n\n## Section 1\n\nBody one.\n\n## Section 2\n\nBody two."
        path = make_md_file(content)
        doc = parse_document(path)
        assert doc.page_count >= 2
        headings = [p["metadata"].get("heading") for p in doc.pages]
        assert "Section 1" in headings
        assert "Section 2" in headings
        os.unlink(path)

    def test_markdown_syntax_stripped(self):
        path = make_md_file("# Title\n\n**bold text** and *italic* and [link](http://example.com).")
        doc = parse_document(path)
        assert "**" not in doc.full_text
        assert "*italic*" not in doc.full_text
        assert "http://example.com" not in doc.full_text
        assert "bold text" in doc.full_text
        assert "italic" in doc.full_text
        assert "link" in doc.full_text
        os.unlink(path)

    def test_plain_txt_file(self):
        path = make_md_file("Plain text content.\nNo markdown here.", suffix=".txt")
        doc = parse_document(path)
        assert doc.file_type == "txt"
        assert "Plain text content" in doc.full_text
        os.unlink(path)

    def test_word_count(self):
        path = make_md_file("one two three four five")
        doc = parse_document(path)
        assert doc.word_count == 5
        os.unlink(path)


# ── PDF tests ─────────────────────────────────────────────────────────────────

class TestPDFParser:
    def test_basic_parse(self):
        path = make_pdf_file()
        doc = parse_document(path)
        assert doc.file_type == "pdf"
        assert doc.page_count == 2
        assert len(doc.full_text) > 50
        os.unlink(path)

    def test_text_content(self):
        path = make_pdf_file()
        doc = parse_document(path)
        assert "RAG" in doc.full_text or "Retrieval" in doc.full_text
        os.unlink(path)

    def test_page_structure(self):
        path = make_pdf_file()
        doc = parse_document(path)
        for page in doc.pages:
            assert "page_num" in page
            assert "text" in page
            assert "metadata" in page
        os.unlink(path)

    def test_metadata_populated(self):
        path = make_pdf_file()
        doc = parse_document(path)
        assert "page_count" in doc.metadata
        assert doc.metadata["page_count"] == 2
        assert "file_size_bytes" in doc.metadata
        os.unlink(path)


# ── DOCX tests ────────────────────────────────────────────────────────────────

class TestDOCXParser:
    def test_basic_parse(self):
        path = make_docx_file()
        doc = parse_document(path)
        assert doc.file_type == "docx"
        assert "Research Summary" in doc.full_text or "Methodology" in doc.full_text
        os.unlink(path)

    def test_headings_become_sections(self):
        path = make_docx_file()
        doc = parse_document(path)
        headings = [p["metadata"].get("heading") for p in doc.pages if p["metadata"].get("heading")]
        assert len(headings) >= 1
        os.unlink(path)

    def test_table_extracted(self):
        path = make_docx_file()
        doc = parse_document(path)
        # Table text should contain the cell values
        assert "BM25" in doc.full_text
        assert "0.72" in doc.full_text
        os.unlink(path)


# ── PPTX tests ────────────────────────────────────────────────────────────────

class TestPPTXParser:
    def test_basic_parse(self):
        path = make_pptx_file()
        doc = parse_document(path)
        assert doc.file_type == "pptx"
        assert doc.page_count == 2
        os.unlink(path)

    def test_slide_titles_extracted(self):
        path = make_pptx_file()
        doc = parse_document(path)
        titles = [p["metadata"]["title"] for p in doc.pages]
        assert "RAG Architecture" in titles
        assert "Results" in titles
        os.unlink(path)

    def test_bullet_text_extracted(self):
        path = make_pptx_file()
        doc = parse_document(path)
        assert "Ingestion" in doc.full_text
        assert "Retrieval" in doc.full_text
        os.unlink(path)

    def test_speaker_notes_extracted(self):
        path = make_pptx_file()
        doc = parse_document(path)
        # Notes should be in the full text
        assert "bottleneck" in doc.full_text
        os.unlink(path)


# ── Router tests ──────────────────────────────────────────────────────────────

class TestRouter:
    def test_unsupported_extension_raises(self):
        with pytest.raises(ValueError, match="No parser available"):
            parse_document(Path("document.xyz"))

    def test_missing_file_raises(self):
        with pytest.raises(FileNotFoundError):
            parse_document(Path("/tmp/does_not_exist_12345.pdf"))

    def test_supported_extensions_list(self):
        exts = get_supported_extensions()
        assert ".pdf" in exts
        assert ".docx" in exts
        assert ".pptx" in exts
        assert ".md" in exts

    def test_parsed_document_interface(self):
        """All parsers return objects with the same interface."""
        files = [make_md_file("# Test\n\nContent."), make_pdf_file(), make_docx_file(), make_pptx_file()]
        for f in files:
            doc = parse_document(f)
            assert hasattr(doc, "full_text")
            assert hasattr(doc, "pages")
            assert hasattr(doc, "metadata")
            assert hasattr(doc, "file_type")
            assert hasattr(doc, "title")
            assert isinstance(doc.full_text, str)
            assert isinstance(doc.pages, list)
            assert len(doc.full_text) > 0
            os.unlink(f)


if __name__ == "__main__":
    pytest.main([__file__, "-v"])